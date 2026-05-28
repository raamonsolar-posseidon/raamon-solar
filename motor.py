"""
Motor de generación de propuestas RAAMON SOLAR
Genera imágenes financieras e inserta en el PPTX
"""
import os, shutil, math
from pathlib import Path
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.ticker as mticker
import numpy as np
from pptx import Presentation

BASE    = Path(__file__).parent
NS_A    = 'http://schemas.openxmlformats.org/drawingml/2006/main'
NS_R    = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
DPI     = 180

DARK_BG   = '#0c0b0b'
DARK_DEEP = '#0a0909'
DARK_CARD = '#141212'
GOLD      = '#F5C518'
GOLD2     = '#c9a00f'
BLUE_BAR  = '#1e55a0'
WHITE     = '#f0ede8'
MUTED     = '#9a9590'

def fmt(n):
    return f"$ {int(round(n)):,}".replace(",", ".")

# ── CÁLCULOS FINANCIEROS ──────────────────────────────────────
def calcular(inversion, paneles, gen_mensual,
             tarifa=1000, inflacion=0.08, degradacion=0.005,
             mant_panel=25000, años=20):
    mant_anual = paneles * mant_panel * 2
    flujos, acums, acum = [], [], -inversion
    for i in range(1, años+1):
        g = gen_mensual * 12 * (1-degradacion)**(i-1)
        t = tarifa * (1+inflacion)**(i-1)
        f = g*t - mant_anual
        acum += f
        flujos.append(f); acums.append(acum)
    acum_tmp, payback = -inversion, años
    for i,f in enumerate(flujos):
        prev = acum_tmp; acum_tmp += f
        if acum_tmp >= 0:
            payback = i + (-prev/f); break
    ahorro20 = sum(flujos)
    van = -inversion + sum(f/(1.12)**i for i,f in enumerate(flujos,1))
    lo, hi = 0.0, 5.0
    for _ in range(200):
        m=(lo+hi)/2
        npv=-inversion+sum(f/(1+m)**i for i,f in enumerate(flujos,1))
        if npv>0: lo=m
        else: hi=m
    tir=(lo+hi)/2
    bc=sum(f/(1.12)**i for i,f in enumerate(flujos,1))/inversion
    return dict(flujos=flujos,acums=acums,payback=payback,
                ahorro1=flujos[0],ahorro20=ahorro20,
                van=van,tir=tir*100,bc=bc,
                roi20=ahorro20/inversion*100,
                roi_anual=ahorro20/inversion*100/años,
                mant_anual=mant_anual)

# ── SLIDE 9 ─ Resumen ejecutivo ───────────────────────────────
def gen_s9(d, paneles, gen_mes, gen_anual, inversion, path):
    W,H = 12.179, 11.025
    fig,ax = plt.subplots(figsize=(W,H))
    fig.patch.set_facecolor(DARK_BG); ax.set_facecolor(DARK_BG); ax.axis('off')
    ax.plot([0.055,0.20],[0.965,0.965],color=GOLD,lw=5,transform=ax.transAxes,solid_capstyle='round')
    ax.text(0.055,0.935,'ANÁLISIS FINANCIERO',fontsize=34,fontweight='bold',color=GOLD,transform=ax.transAxes,va='top')
    ax.text(0.055,0.874,'SISTEMA FOTOVOLTAICO',fontsize=42,fontweight='bold',color=WHITE,transform=ax.transAxes,va='top')
    ax.text(0.055,0.820,'INVERSIÓN INTELIGENTE · AHORRO GARANTIZADO',fontsize=14,color=MUTED,transform=ax.transAxes,va='top',fontstyle='italic')
    ax.text(0.920,0.955,'⚡ RA-AMON\n      SOLAR',fontsize=15,fontweight='bold',color=DARK_BG,
            transform=ax.transAxes,ha='center',va='top',
            bbox=dict(boxstyle='round,pad=0.5',facecolor='#1a1818',edgecolor='#3a3838',lw=1.5))
    ax.plot([0.055,0.945],[0.805,0.805],color=GOLD2,lw=1,alpha=0.35,transform=ax.transAxes)
    ax.text(0.055,0.787,'RESUMEN EJECUTIVO',fontsize=16,fontweight='bold',color=GOLD,transform=ax.transAxes,va='top')

    def card(x,y,w,h,title,value,sub):
        ax.add_patch(mpatches.FancyBboxPatch((x,y),w,h,boxstyle='round,pad=0.012',
            lw=2,edgecolor=GOLD2,facecolor=DARK_CARD,transform=ax.transAxes,zorder=2))
        ax.text(x+w/2,y+h-0.025,title,fontsize=11.5,color=MUTED,transform=ax.transAxes,
                ha='center',va='top',fontweight='bold',multialignment='center')
        ax.text(x+w/2,y+h/2,value,fontsize=16,color=WHITE,transform=ax.transAxes,
                ha='center',va='center',fontweight='bold')
        if sub: ax.text(x+w/2,y+0.022,sub,fontsize=10,color=GOLD2,transform=ax.transAxes,ha='center',va='bottom')

    gap=0.012; cw=(0.890-gap*3)/4; ch=0.190
    r1=[('INVERSIÓN INICIAL',fmt(inversion),'COP'),
        ('AHORRO NETO  AÑO 1',fmt(d['ahorro1']),'Año 1'),
        ('RECUPERACIÓN INVERSIÓN',f"~ {d['payback']:.1f} AÑOS",''),
        ('ROI  20 AÑOS',f"~ {d['roi20']:,.0f} %".replace(',','.'),'')]
    r2=[('GENERACIÓN MENSUAL',f'{gen_mes:,} kWh/mes'.replace(',','.'),''),
        ('GENERACIÓN ANUAL',f'{gen_anual:,} kWh/año'.replace(',','.'),''),
        ('AHORRO ACUMULADO 20 AÑOS',fmt(d['ahorro20']),'COP'),
        ('VIDA ÚTIL DEL SISTEMA','+ 25 AÑOS','')]
    for i,(t,v,s) in enumerate(r1): card(0.055+i*(cw+gap),0.558,cw,ch,t,v,s)
    for i,(t,v,s) in enumerate(r2): card(0.055+i*(cw+gap),0.345,cw,ch,t,v,s)
    ax.plot([0.055,0.945],[0.333,0.333],color=GOLD2,lw=1,alpha=0.25,transform=ax.transAxes)
    ax.text(0.055,0.314,'TEXTO EJECUTIVO',fontsize=14,fontweight='bold',color=GOLD,transform=ax.transAxes,va='top')
    lines=[f"El sistema de {paneles} paneles AE Solar 700W recupera la inversión de {fmt(inversion)} en aproximadamente {d['payback']:.1f} años.",
           f"Genera ahorros acumulados superiores a {fmt(d['ahorro20'])} COP durante 20 años de operación proyectada.",
           f"TIR {d['tir']:.0f}%  ·  VAN {fmt(d['van'])} COP  ·  Relación B/C {d['bc']:.1f}:1  ·  ROI anual {d['roi_anual']:.1f}%"]
    for i,ln in enumerate(lines):
        ax.text(0.055,0.276-i*0.068,ln,fontsize=12.5,color=MUTED,transform=ax.transAxes,va='top')
    plt.subplots_adjust(left=0,right=1,top=1,bottom=0)
    plt.savefig(path,dpi=DPI,facecolor=DARK_BG,edgecolor='none'); plt.close()

# ── SLIDE 10 ─ Gráfica proyección ────────────────────────────
def gen_s10(d, años_list, inversion, path):
    W,H=16.753,10.780; fig=plt.figure(figsize=(W,H)); fig.patch.set_facecolor(DARK_BG)
    ax_l=fig.add_axes([0.00,0.30,0.13,0.45]); ax_l.set_facecolor(DARK_BG); ax_l.axis('off')
    ax_l.text(0.05,0.70,'Analisis',fontsize=22,color=GOLD,fontweight='bold',transform=ax_l.transAxes)
    ax_l.text(0.05,0.28,'Financiero',fontsize=22,color=WHITE,fontweight='bold',transform=ax_l.transAxes)
    ax=fig.add_axes([0.215,0.24,0.765,0.72]); ax.set_facecolor('#080707')
    ya=np.array(años_list); fa=np.array(d['flujos']); aa=np.array(d['acums'])
    ax.bar(ya,fa,color=BLUE_BAR,alpha=0.88,width=0.65,zorder=2)
    ax2=ax.twinx(); ax2.set_facecolor('#080707')
    ax2.plot(ya,aa,color=GOLD,lw=3,marker='o',ms=5,zorder=3,markerfacecolor=GOLD,markeredgecolor=DARK_BG,markeredgewidth=1.5)
    ax2.axhline(0,color=MUTED,lw=1,ls='--',alpha=0.4)
    pb_x=años_list[0]+d['payback']
    ax.axvline(pb_x,color=GOLD,lw=2,ls='--',alpha=0.8,zorder=4)
    yr1=int(años_list[0]+d['payback']); yr2=yr1+1
    ax.annotate(f'PUNTO DE EQUILIBRIO\nALCANZADO\nEntre {yr1} y {yr2}',
                xy=(pb_x,0),xytext=(pb_x+2.5,max(fa)*0.38),fontsize=9,color=DARK_BG,fontweight='bold',
                bbox=dict(boxstyle='round,pad=0.5',facecolor=GOLD,edgecolor=GOLD2),
                arrowprops=dict(arrowstyle='->',color=GOLD,lw=2),zorder=5)
    ax2.annotate(f'-{fmt(inversion)}',xy=(años_list[0],-inversion),xytext=(años_list[0]+0.5,-inversion*1.5),
                 fontsize=8,color=WHITE,bbox=dict(boxstyle='round,pad=0.3',facecolor='#252323',edgecolor=MUTED))
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x,p: f'$ {int(x/1e6)}.000.000' if x>=0 else f'-$ {int(abs(x)/1e6)}.000.000'))
    ax2.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x,p: f'$ {x/1e9:.2f}B' if abs(x)>=1e9 else f'$ {int(x/1e6)}M'))
    ax.tick_params(colors=MUTED,labelsize=8); ax2.tick_params(colors=MUTED,labelsize=8)
    ax.set_xticks(ya); ax.set_xticklabels([str(y) for y in años_list],rotation=45,fontsize=8,color=MUTED)
    for sp in ['bottom','left']: ax.spines[sp].set_color('#3a3838')
    for sp in ['top','right']: ax.spines[sp].set_visible(False)
    ax2.spines['right'].set_color('#3a3838'); ax2.spines['top'].set_visible(False)
    ax.grid(axis='y',color='#2a2828',alpha=0.5,ls='--',lw=0.8)
    ax.text(-0.01,1.05,'PROYECCIÓN DE AHORRO ACUMULADO',fontsize=14,fontweight='bold',color=WHITE,transform=ax.transAxes,va='bottom')
    p1=mpatches.Patch(facecolor=BLUE_BAR,label='Ahorro anual neto')
    p2=plt.Line2D([0],[0],color=GOLD,lw=3,label='Ahorro acumulado')
    ax.legend(handles=[p1,p2],loc='upper left',fontsize=9,facecolor=DARK_CARD,labelcolor=WHITE,edgecolor='#3a3838',framealpha=0.95,bbox_to_anchor=(0.0,0.99))
    ax.set_ylabel('AHORRO ANUAL (COP)',color=MUTED,fontsize=8,labelpad=4)
    ax2.set_ylabel('AHORRO ACUMULADO (COP)',color=MUTED,fontsize=8,labelpad=4,rotation=270)
    ax.set_xlabel('AÑO',color=MUTED,fontsize=8,labelpad=4)
    ax_i=fig.add_axes([0.215,0.01,0.765,0.21]); ax_i.set_facecolor(DARK_BG); ax_i.axis('off')
    info=[('La inversión se recupera en',f"aproximadamente {d['payback']:.1f} años",'gracias al ahorro generado.'),
          ('A partir del año 3 el sistema','genera ahorros netos crecientes','durante más de 20 años.'),
          ('Ahorro acumulado proyectado','a 20 años de más de',fmt(d['ahorro20'])+' COP')]
    for i,(l1,l2,l3) in enumerate(info):
        bx=0.01+i*0.335
        ax_i.add_patch(mpatches.FancyBboxPatch((bx,0.05),0.315,0.90,boxstyle='round,pad=0.02',lw=1.5,edgecolor=GOLD2,facecolor=DARK_CARD,transform=ax_i.transAxes))
        ax_i.text(bx+0.158,0.85,l1,fontsize=10,color=WHITE,fontweight='bold',transform=ax_i.transAxes,ha='center',va='top')
        ax_i.text(bx+0.158,0.54,l2,fontsize=10,color=MUTED,transform=ax_i.transAxes,ha='center',va='top')
        col=GOLD if i==2 else MUTED
        ax_i.text(bx+0.158,0.23,l3,fontsize=10,color=col,fontweight='bold' if i==2 else 'normal',transform=ax_i.transAxes,ha='center',va='top')
    plt.savefig(path,dpi=DPI,facecolor=DARK_BG,edgecolor='none',bbox_inches='tight'); plt.close()

# ── SLIDE 11 ─ Comparativa ────────────────────────────────────
def gen_s11(d, inversion, paneles, gen_mes, path):
    costo_acum=sum(gen_mes*12*1000*(1.08**i) for i in range(20))
    W,H=15.813,8.799; fig,ax=plt.subplots(figsize=(W,H))
    fig.patch.set_facecolor(DARK_BG); ax.set_facecolor(DARK_BG); ax.axis('off')
    ax.text(0.00,0.98,'Analisis ',fontsize=38,color=GOLD,fontweight='bold',transform=ax.transAxes,va='top')
    ax.text(0.26,0.98,'Financiero',fontsize=38,color=WHITE,fontweight='bold',transform=ax.transAxes,va='top')
    ax.add_patch(mpatches.FancyBboxPatch((0.02,0.06),0.96,0.83,boxstyle='round,pad=0.012',lw=2,edgecolor=GOLD2,facecolor='#0c0b0b',transform=ax.transAxes))
    ax.text(0.50,0.875,'COMPARATIVA: SEGUIR PAGANDO vs INVERTIR EN SOLAR',fontsize=14,fontweight='bold',color=GOLD,transform=ax.transAxes,ha='center')
    ax.plot([0.50,0.50],[0.12,0.84],color='#3a3838',lw=1.5,transform=ax.transAxes)
    ax.add_patch(plt.Circle((0.50,0.48),0.042,color='#1a1818',transform=ax.transAxes,zorder=5,clip_on=False))
    ax.text(0.50,0.48,'VS',fontsize=15,fontweight='bold',color=GOLD,transform=ax.transAxes,ha='center',va='center',zorder=6)
    ax.text(0.245,0.820,'✗   SEGUIR PAGANDO ENERGÍA',fontsize=14,fontweight='bold',color='#ff5555',transform=ax.transAxes,ha='center')
    for i,item in enumerate([f'Pago acumulado 20 años: {fmt(costo_acum)} COP','Costo energético creciente cada año','Dinero que se pierde y no genera activo','Expuesto a la inflación energética']):
        ax.text(0.06,0.730-i*0.140,f'◆   {item}',fontsize=12,color=WHITE,transform=ax.transAxes)
    ax.text(0.755,0.820,'✓   INVERTIR EN ENERGÍA SOLAR',fontsize=14,fontweight='bold',color='#44cc44',transform=ax.transAxes,ha='center')
    for i,item in enumerate([f'Inversión inicial: {fmt(inversion)} COP',f"Ahorro acumulado 20 años: {fmt(d['ahorro20'])} COP",'Activos que aumentan el valor del inmueble','Protección contra alzas tarifarias']):
        ax.text(0.535,0.730-i*0.140,f'◆   {item}',fontsize=12,color=WHITE,transform=ax.transAxes)
    ax.add_patch(mpatches.FancyBboxPatch((0.06,0.08),0.88,0.190,boxstyle='round,pad=0.012',lw=2,edgecolor=GOLD,facecolor='#131110',transform=ax.transAxes))
    ax.text(0.165,0.175,'☀',fontsize=36,color=GOLD,transform=ax.transAxes,va='center')
    ax.text(0.248,0.208,'Invertir en solar es transformar un gasto creciente',fontsize=13,color=WHITE,transform=ax.transAxes)
    ax.text(0.248,0.132,'en un ',fontsize=13,color=WHITE,transform=ax.transAxes)
    ax.text(0.328,0.132,'activo rentable y sostenible.',fontsize=13,fontweight='bold',color=GOLD,transform=ax.transAxes)
    plt.subplots_adjust(left=0,right=1,top=1,bottom=0)
    plt.savefig(path,dpi=DPI,facecolor=DARK_BG,edgecolor='none'); plt.close()

# ── SLIDE 12 ─ Indicadores ────────────────────────────────────
def gen_s12(d, path):
    X0=0.0326; W,H=20.6743,9.5719
    fig,ax=plt.subplots(figsize=(W,H))
    fig.patch.set_facecolor(DARK_BG); ax.set_facecolor(DARK_BG); ax.axis('off')
    ax.text(X0+0.010,0.958,'INDICADORES FINANCIEROS CLAVE',fontsize=32,fontweight='bold',color=GOLD,transform=ax.transAxes,va='top',zorder=10)
    ax.plot([X0,0.998],[0.862,0.862],color=GOLD2,lw=1.5,alpha=0.4,transform=ax.transAxes,zorder=8)
    vw=1.0-X0; cg=0.008; cw=(vw-cg*3)/4; ch=0.790; cy=0.058
    inds=[('VAN\n(Valor Actual Neto)',fmt(d['van']),'(> 0)\nProyecto viable'),
          ('TIR\n(Tasa Interna de Retorno)',f"~ {d['tir']:.0f} %",'Altamente rentable'),
          ('Relación B/C\n(Beneficio / Costo)',f"~ {d['bc']:.1f} : 1",f"Cada $1 invertido\ngenera ${d['bc']:.1f}"),
          ('ROI Anual Promedio',f"~ {d['roi_anual']:.1f} %",'Promedio durante\n20 años')]
    syms=['◈','◉','◎','◇']
    for i,(title,value,sub) in enumerate(inds):
        x=X0+i*(cw+cg)
        ax.add_patch(mpatches.FancyBboxPatch((x,cy),cw,ch,boxstyle='round,pad=0.008',lw=2.5,edgecolor=GOLD2,facecolor=DARK_DEEP,transform=ax.transAxes,zorder=2))
        ax.text(x+cw/2,cy+ch-0.025,title,fontsize=17,color=MUTED,transform=ax.transAxes,ha='center',va='top',fontweight='bold',multialignment='center')
        iy=cy+ch*0.58
        ax.add_patch(mpatches.Ellipse((x+cw/2,iy),cw*0.52,0.180,facecolor=GOLD2,alpha=0.55,transform=ax.transAxes,zorder=3))
        ax.text(x+cw/2,iy,syms[i],fontsize=46,color=GOLD,transform=ax.transAxes,ha='center',va='center',fontweight='bold',zorder=4)
        ax.text(x+cw/2,cy+ch*0.32,value,fontsize=30,color=WHITE,transform=ax.transAxes,ha='center',va='center',fontweight='bold')
        ax.text(x+cw/2,cy+0.025,sub,fontsize=14,color=MUTED,transform=ax.transAxes,ha='center',va='bottom',multialignment='center')
    plt.subplots_adjust(left=0,right=1,top=1,bottom=0)
    plt.savefig(path,dpi=DPI,facecolor=DARK_BG,edgecolor='none'); plt.close()

# ── SLIDE 19 ─ Inversor ───────────────────────────────────────
def gen_s19(marca, modelo, kw, fase, path):
    W,H=13.935,9.296; rng=np.random.default_rng(42)
    fig,ax=plt.subplots(figsize=(W,H))
    fig.patch.set_facecolor(DARK_BG); ax.set_facecolor(DARK_BG); ax.axis('off')
    ax.scatter(rng.uniform(0,1,280),rng.uniform(0,1,280),s=rng.uniform(0.3,2.0,280),c='white',alpha=0.13,transform=ax.transAxes,zorder=1)
    for r,a in [(0.50,0.03),(0.36,0.06),(0.22,0.10),(0.12,0.14)]:
        ax.add_patch(mpatches.Ellipse((0.72,0.50),r*2,r*2.2,facecolor='#0d1f3c',alpha=a,transform=ax.transAxes,zorder=1))
    ax.text(0.938,0.965,'#1 EN INVERSORES\nSOLARES MUNDIAL',fontsize=13,fontweight='bold',color=DARK_BG,
            transform=ax.transAxes,ha='center',va='top',zorder=10,
            bbox=dict(boxstyle='round,pad=0.45',facecolor=GOLD,edgecolor=GOLD2,lw=2))
    ax.text(0.050,0.820,'INVERSOR ON GRID',fontsize=20,fontweight='bold',color=WHITE,transform=ax.transAxes,va='top',zorder=10)
    ax.text(0.050,0.772,f'{marca.upper()} {modelo.upper()}',fontsize=36,fontweight='bold',color=GOLD,transform=ax.transAxes,va='top',zorder=10)
    ax.text(0.050,0.706,f'{kw}kW  {fase.upper()}',fontsize=28,fontweight='bold',color=WHITE,transform=ax.transAxes,va='top',zorder=10)
    ax.text(0.050,0.658,f'⚡  {fase.upper()}  —  On Grid',fontsize=13,fontweight='bold',color=GOLD,transform=ax.transAxes,va='top',zorder=10)
    ax.text(0.050,0.622,'Energía inteligente · máxima eficiencia · control total',fontsize=11,color=MUTED,transform=ax.transAxes,va='top',fontstyle='italic',zorder=10)
    feats=[('◈','Alta eficiencia hasta 98,6%','Aprovecha al máximo cada rayo de sol'),
           ('◉','Monitoreo inteligente por Wi-Fi','Controla tu sistema desde el celular'),
           ('◎','2 MPPT independientes','Mayor producción con sombras parciales'),
           ('◈','Protección total con IA','Detección de arco, anti-isla integrada'),
           ('◉',f'Potencia nominal: {kw} kW',f'Inversor {fase} certificado')]
    for i,(icon,title,desc) in enumerate(feats):
        y=0.565-i*0.110
        ax.add_patch(plt.Circle((0.060,y+0.004),0.018,facecolor=GOLD2+'22',edgecolor=GOLD2,lw=1.3,transform=ax.transAxes,zorder=9,clip_on=False))
        ax.text(0.060,y+0.004,icon,fontsize=11,color=GOLD,transform=ax.transAxes,ha='center',va='center',zorder=10)
        ax.text(0.090,y+0.018,title,fontsize=12,fontweight='bold',color=WHITE,transform=ax.transAxes,va='top',zorder=10)
        ax.text(0.090,y-0.016,desc,fontsize=10,color=MUTED,transform=ax.transAxes,va='top',zorder=10)
    ax.add_patch(mpatches.FancyBboxPatch((0.510,0.115),0.340,0.755,boxstyle='round,pad=0.015',lw=2.5,edgecolor='#c8c4c0',facecolor='#eae7e2',transform=ax.transAxes,zorder=8))
    ax.text(0.680,0.808,marca.upper(),fontsize=20,fontweight='bold',color='#cc0000',transform=ax.transAxes,ha='center',va='center',zorder=10)
    ax.text(0.680,0.763,modelo,fontsize=9,color='#666',transform=ax.transAxes,ha='center',va='center',zorder=10)
    ax.add_patch(mpatches.FancyBboxPatch((0.530,0.570),0.128,0.088,boxstyle='round,pad=0.008',lw=1.5,edgecolor='#555',facecolor='#121f12',transform=ax.transAxes,zorder=9))
    ax.text(0.594,0.615,f'{kw}.00 kW',fontsize=11,color='#66ee66',transform=ax.transAxes,ha='center',va='center',zorder=10,fontfamily='monospace')
    ax.add_patch(mpatches.FancyBboxPatch((0.528,0.370),0.210,0.130,boxstyle='round,pad=0.010',lw=2,edgecolor=GOLD,facecolor='#0e0d0d',transform=ax.transAxes,zorder=10))
    ax.text(0.633,0.448,fase.upper(),fontsize=14,fontweight='bold',color=GOLD,transform=ax.transAxes,ha='center',va='center',zorder=11)
    ax.text(0.633,0.392,'220V / 400V',fontsize=10,color=WHITE,transform=ax.transAxes,ha='center',va='center',zorder=11)
    for xi in np.linspace(0.530,0.835,8):
        ax.add_patch(plt.Rectangle((xi,0.124),0.020,0.050,facecolor='#888',edgecolor='#666',lw=0.5,transform=ax.transAxes,zorder=9))
    gx,gy=0.905,0.590
    ax.add_patch(mpatches.FancyBboxPatch((gx-0.068,gy-0.165),0.136,0.245,boxstyle='round,pad=0.015',lw=2.5,edgecolor=GOLD,facecolor='#1a1300',transform=ax.transAxes,zorder=10))
    for r,a in [(0.068,0.40),(0.050,0.55),(0.033,0.70)]:
        ax.add_patch(plt.Circle((gx,gy+0.015),r,facecolor='none',edgecolor=GOLD,lw=1.6,alpha=a,transform=ax.transAxes,zorder=11))
    ax.text(gx,gy+0.018,'10',fontsize=30,fontweight='bold',color=GOLD,transform=ax.transAxes,ha='center',va='center',zorder=12)
    ax.text(gx,gy-0.068,'AÑOS',fontsize=11,fontweight='bold',color=WHITE,transform=ax.transAxes,ha='center',va='center',zorder=12)
    ax.text(gx,gy-0.104,'GARANTÍA',fontsize=9.5,color=MUTED,transform=ax.transAxes,ha='center',va='center',zorder=12)
    ax.plot([0.028,0.972],[0.100,0.100],color=GOLD2,lw=1.5,alpha=0.5,transform=ax.transAxes,zorder=8)
    ax.text(0.500,0.065,'Convierte el sol en',fontsize=13,color=WHITE,transform=ax.transAxes,ha='center',va='center',zorder=10)
    ax.text(0.500,0.028,f'AHORRO TOTAL  con  {marca.upper()}',fontsize=20,fontweight='bold',color=GOLD,transform=ax.transAxes,ha='center',va='center',zorder=10)
    plt.subplots_adjust(left=0,right=1,top=1,bottom=0)
    plt.savefig(path,dpi=DPI,facecolor=DARK_BG,edgecolor='none'); plt.close()

# ── ENSAMBLAR PPTX ────────────────────────────────────────────
def reemplazar_blip(slide, old_rid, new_rid):
    for shape in slide.shapes:
        for blip in shape._element.iter(f'{{{NS_A}}}blip'):
            if blip.get(f'{{{NS_R}}}embed') == old_rid:
                blip.set(f'{{{NS_R}}}embed', new_rid)

def limpiar_fill(slide, shape_name):
    for shape in slide.shapes:
        if shape.name == shape_name:
            for fr in shape._element.iter(f'{{{NS_A}}}fillRect'):
                fr.set('l','0'); fr.set('t','0'); fr.set('r','0'); fr.set('b','0')

def generar_propuesta(nombre, paneles, precio, inv_marca, inv_modelo, inv_kw, inv_fase, tmp_dir):
    plantilla = BASE / 'propuesta_plantilla.pptx'
    output    = Path(tmp_dir) / f"PROPUESTA_{nombre.replace(' ','_').upper()}.pptx"

    cap_kw    = round(paneles * 0.7, 1)
    gen_mes   = paneles * 84
    gen_anual = gen_mes * 12
    años_list = list(range(2025, 2045))
    d = calcular(precio, paneles, gen_mes)

    # Generar imágenes
    imgs = {
        's9' : tmp_dir+'/s9.png',
        's10': tmp_dir+'/s10.png',
        's11': tmp_dir+'/s11.png',
        's12': tmp_dir+'/s12.png',
        's19': tmp_dir+'/s19.png',
    }
    gen_s9 (d, paneles, gen_mes, gen_anual, precio, imgs['s9'])
    gen_s10(d, años_list, precio, imgs['s10'])
    gen_s11(d, precio, paneles, gen_mes, imgs['s11'])
    gen_s12(d, imgs['s12'])
    gen_s19(inv_marca, inv_modelo, inv_kw, inv_fase, imgs['s19'])

    # Ensamblar
    shutil.copy2(plantilla, output)
    prs = Presentation(str(output))

    # Slide 2 — nombre
    for shape in prs.slides[1].shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if 'ING. OMAR ORTIZ' in run.text:
                        run.text = nombre.upper()

    # Slide 3 — inversor en listado
    inv_label = f'Inversor {inv_marca.upper()} {inv_modelo.upper()} {inv_kw}KW {inv_fase.upper()}'
    for shape in prs.slides[2].shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if 'SOLAX 10KW TRIFASICO' in run.text:
                        run.text = run.text.replace('Inversor SOLAX 10KW TRIFASICO', inv_label)

    # Slide 6 — datos proyecto
    cap_str = f'    {int(cap_kw)}  ' if cap_kw==int(cap_kw) else f'    {cap_kw}  '
    precio_fmt = f"$ {precio:,}".replace(',','.')
    for shape in prs.slides[5].shapes:
        if not shape.has_text_frame: continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                t=run.text
                if t=='    14  ':               run.text=cap_str
                elif t=='1680 ':               run.text=f'{gen_mes} '
                elif t=='$ 43.500.000':        run.text=precio_fmt
                elif t=='20 ':                 run.text=f'{paneles} '
                elif t=='SOLAX 10KW TRIFASICO': run.text=f'{inv_marca.upper()} {inv_modelo.upper()} {inv_kw}KW {inv_fase.upper()}'

    # Slides 9-12 financieras
    for si,sn,rid,img in [(8,'Freeform 2','rId2',imgs['s9']),
                           (9,'Freeform 2','rId2',imgs['s10']),
                           (10,'Freeform 3','rId3',imgs['s11'])]:
        sl=prs.slides[si]; _,nr=sl.part.get_or_add_image_part(img)
        reemplazar_blip(sl,rid,nr); limpiar_fill(sl,sn)
    sl12=prs.slides[11]; _,nr12=sl12.part.get_or_add_image_part(imgs['s12'])
    reemplazar_blip(sl12,'rId3',nr12); limpiar_fill(sl12,'Freeform 3')

    # Slide 19 inversor
    sl19=prs.slides[18]; _,nr19=sl19.part.get_or_add_image_part(imgs['s19'])
    reemplazar_blip(sl19,'rId4',nr19); limpiar_fill(sl19,'Freeform 5')

    prs.save(str(output))
    return str(output), d
